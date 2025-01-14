from pdfminer.high_level import extract_text
import os
import zipfile
from docx import Document
from pptx import Presentation
import openpyxl
from pdf2image import convert_from_path
from PIL import Image as PILImage
from io import BytesIO
import io
import cv2
from tqdm import tqdm
import numpy as np
import pytesseract
from PIL import Image
pytesseract.pytesseract.tesseract_cmd =r'/usr/bin/tesseract'

# Functions to set up OCR
def deskew_pptx(image):
    if image.mode != 'L':  # 'L' stands for grayscale mode
        image = image.convert('L')
    return image

def convert_wmf_to_png(image_bytes):
    with Image.open(io.BytesIO(image_bytes)) as img:
        with io.BytesIO() as output:
            img.save(output, format="PNG")
            return output.getvalue()

def deskew(image):
    if len(image.shape) == 2:  # Grayscale image (height, width)
        gray = image
    else:  # Color image (height, width, channels)
        gray = cv2.cvtColor(image, cv2.COLOR_BGR2GRAY)
    
    gray = cv2.bitwise_not(gray)
    coords = np.column_stack(np.where(gray > 0))
    angle = cv2.minAreaRect(coords)[-1]
    
    if angle < -45:
        angle = -(90 + angle)
    else:
        angle = -angle

    (h, w) = image.shape[:2]
    center = (w // 2, h // 2)
    M = cv2.getRotationMatrix2D(center, angle, 1.0)
    rotated = cv2.warpAffine(image, M, (w, h), flags=cv2.INTER_CUBIC, borderMode=cv2.BORDER_REPLICATE)

    return rotated

def extract_text_from_image(image):
    text = pytesseract.image_to_string(image)
    return text

def extract_images_from_docx(docx_path):
    images = []
    with zipfile.ZipFile(docx_path, 'r') as docx_zip:
        # Check for images in the 'word/media/' folder
        for file_name in docx_zip.namelist():
            if file_name.startswith('word/media/'):
                image_data = docx_zip.read(file_name)
                images.append(BytesIO(image_data))  # Store image in memory
    return images

def ocr_from_images(images):
    text = ""
    for img_data in images:
        img = Image.open(img_data)
        text += pytesseract.image_to_string(img)
        text += "\n"
    return text

def get_file_text(extension, file):
    if extension == 'pdf':
        return get_one_pdf(file)
    if extension == 'docx':
        return get_one_docx(file)
    if extension == 'pptx':
        return get_one_pptx(file)
    if extension == 'xlsx':
        return get_one_xlsx(file)


def get_one_pdf(pdf):
    pages = convert_from_path(pdf)
    for page in pages:
        preprocessed_image = deskew(np.array(page))
        return extract_text_from_image(preprocessed_image)
        
def get_one_docx(docx):
    doc = Document(docx)
    
    for para in doc.paragraphs:
        text += para.text + "\n"
    
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                text += cell.text + " "
            text += "\n"
    
    images = extract_images_from_docx(docx)
    text += ocr_from_images(images)
    return text

# Function to extract text from PDF files using pdfminer.six
def get_pdf_text(pdf_docs):
    print("Get pdf text")
    extracted_text = ""
    for pdf in tqdm(pdf_docs):
        try:
            extracted_text += get_one_pdf(pdf)
        except:
            print(f'failed to get {pdf}, skipping')
    return extracted_text

# Function to extract text from DOCX files including text from images using OCR
def get_docx_text(docx_docs):
    print("Get docx text")
    text = ""
    
    for docx in tqdm(docx_docs):
        try:
            text += get_one_docx(docx)
        except:
            print(f'failed to get {docx}, skipping')
    return text

def get_one_pptx(pptx):
    print(f"Processing PPTX document: {pptx}")
    text = ''
    prs = Presentation(pptx)
    
    for slide_index, slide in enumerate(prs.slides):
        print(f"Processing slide {slide_index + 1}/{len(prs.slides)}")
        for shape_index, shape in enumerate(slide.shapes):
            print(f"Processing shape {shape_index + 1}/{len(slide.shapes)}")
            if shape.shape_type == 13:  # This is the image shape type
                try:
                    print("Found image shape")
                    # Extract image bytes
                    image = shape.image
                    image_bytes = image.blob
                    
                    # Open image with Pillow
                    try:
                        img = Image.open(io.BytesIO(image_bytes))
                        print("Image opened successfully")
                    except IOError:
                        print("Image is WMF, converting to PNG")
                        # If the image is WMF, convert it to PNG
                        image_bytes = convert_wmf_to_png(image_bytes)
                        img = Image.open(io.BytesIO(image_bytes))
                        print("WMF image converted to PNG and opened successfully")
                    
                    # Handle transparency and convert image to RGBA (if it's a palette-based image with transparency)
                    if img.mode == 'RGBA':
                        img = img.convert('RGBA')  # Ensure it's RGBA if it's not
                    elif img.mode == 'P' and 'transparency' in img.info:
                        img = img.convert('RGBA')  # Convert palette-based image with transparency to RGBA
                    elif img.mode == 'P':
                        img = img.convert('RGB')  # Convert palette image without transparency to RGB
                    elif img.mode == 'LA':
                        img = img.convert('RGBA')  # If image mode is LA (Luminance + Alpha), convert to RGBA
                    elif img.mode == 'L':
                        img = img.convert('RGB')  # Convert grayscale to RGB for OCR
                    
                    print("Image mode converted successfully")
                    
                    # Convert PIL image to NumPy array
                    img_np = np.array(img)
                    print("Image converted to NumPy array")
                    
                    # Apply deskew (if necessary)
                    img_np = deskew(img_np)
                    print("Deskew applied to image")
                    
                    # Perform OCR on the image
                    text += pytesseract.image_to_string(img_np) + "\n"
                    print("OCR performed on image")
                except Exception as e:
                    print(f"Error processing image: {e}")
            
            if hasattr(shape, "text"):
                text += shape.text + "\n"
                print("Text extracted from shape")
            
            # Check if the shape is a table
            if shape.has_table:
                print("Found table shape")
                table = shape.table
                for row_index, row in enumerate(table.rows):
                    for cell_index, cell in enumerate(row.cells):
                        text += cell.text + "\t"  # Add a tab space between table columns
                        print(f"Text extracted from table cell {cell_index + 1}/{len(row.cells)} in row {row_index + 1}/{len(table.rows)}")
                    text += "\n"  # New line after each row
        
        # Extract text from speaker notes (if any)
        if slide.has_notes_slide:
            notes_slide = slide.notes_slide
            if notes_slide:
                notes_text = notes_slide.notes_text_frame.text
                text += f"{notes_text}\n"
                print("Text extracted from speaker notes")
    return text

def get_one_xlsx(xlsx):
    text = ''
    wb = openpyxl.load_workbook(xlsx)
    for sheet in wb.sheetnames:
        ws = wb[sheet]
        for row in ws.iter_rows():
            for cell in row:
                if cell.value is not None:
                    text += str(cell.value) + " "
            text += "\n"
    return text


def get_pptx_text(pptx_docs):
    print("Starting get_pptx_text function")
    text = ""
    
    for pptx in tqdm(pptx_docs):
        try: 
            text += get_one_pptx(pptx)
        except: 
            print(f'failed to get {pptx}, skipping')
    
    print("Finished get_pptx_text function")
    return text

# Function to extract text from XLSX files
def get_xlsx_text(xlsx_docs):
    print("Get xlsx text")
    text = ""
    for xlsx in tqdm(xlsx_docs):
        try:
            text += get_one_xlsx(xlsx)
        except:
            print(f'failed to get {xlsx}, skipping')
    return text

# Directory paths
pdf_dir =  r'modelling/data/pdf_files'
docx_dir = r'modelling/data/docx_files'
pptx_dir = r'modelling/data/pptx_files'
xlsx_dir = r'modelling/data/xlsx_files'

# Fetch file paths
pdf_files = [os.path.join(pdf_dir, f) for f in os.listdir(pdf_dir) if f.endswith('.pdf')]
docx_files = [os.path.join(docx_dir, f) for f in os.listdir(docx_dir) if f.endswith('.docx')]
pptx_files = [os.path.join(pptx_dir, f) for f in os.listdir(pptx_dir) if f.endswith('.pptx')]
xlsx_files = [os.path.join(xlsx_dir, f) for f in os.listdir(xlsx_dir) if f.endswith('.xlsx')]

output_file = os.path.join('modelling', 'extracted_text.txt')

if not os.path.exists(output_file):
    print(f"{output_file} does not exist. Extracting text...")

    # Extract text for each document type
    pdf_text = get_pdf_text(pdf_files)
    docx_text = get_docx_text(docx_files)
    pptx_text = get_pptx_text(pptx_files)
    xlsx_text = get_xlsx_text(xlsx_files)

    # Function to store all extracted text into one file
    def save_all_text_to_file(texts, file_name):
        with open(file_name, 'w', encoding='utf-8') as f:
            for text_type, text in texts.items():
                # f.write(f"--- {text_type.upper()} TEXT ---\n")
                f.write(text + "\n")

    # Combine all extracted texts into a dictionary
    all_texts = {
        "pdf": pdf_text,
        "docx": docx_text,
        "pptx": pptx_text,
        "xlsx": xlsx_text
    }

    # Save all text into one file
    save_all_text_to_file(all_texts, output_file)
    print(f"Extracted text saved to {output_file}")
else:
    print(f"{output_file} already exists. Skipping text extraction.")